import csv
import io
import os
import re
import tempfile

import pandas as pd
import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation
from pypdf import PdfReader
from agno.agent import Agent
from agno.models.google import Gemini
from agno.models.openai import OpenAIChat
from agno.tools.duckdb import DuckDbTools
from agno.tools.pandas import PandasTools

load_dotenv()


def resolve_openai_key(session_state):
    if "openai_key" in session_state and session_state.get("openai_key"):
        return session_state["openai_key"]

    env_key = os.getenv("OPENAI_API_KEY")
    if env_key and env_key.strip():
        session_state["openai_key"] = env_key.strip()
        return session_state["openai_key"]

    return None


def build_model(provider, api_key):
    provider_name = (provider or "google").strip().lower()
    if provider_name == "openai":
        if not api_key:
            raise ValueError("OpenAI API key is required when using the OpenAI provider.")
        return OpenAIChat(id="gpt-4o", api_key=api_key)

    if provider_name == "google" or provider_name == "gemini":
        gemini_api_key = os.getenv("GOOGLE_API_KEY") or api_key
        if not gemini_api_key:
            raise ValueError("Google API key is required when using the Google/Gemini provider.")
        return Gemini(id="gemini-2.0-flash-001", api_key=gemini_api_key)

    raise ValueError(f"Unsupported provider: {provider_name}")


# Function to preprocess and save the uploaded file
def extract_pdf_text(file):
    try:
        reader = PdfReader(file)
        text_chunks = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_chunks.append(page_text.strip())
        return "\n\n".join(text_chunks)
    except Exception as exc:
        st.error(f"Unable to read PDF: {exc}")
        return None


def extract_ppt_text(file):
    try:
        prs = Presentation(file)
        text_chunks = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_chunks.append(f"Slide {idx}:\n" + "\n".join(slide_text))
        return "\n\n".join(text_chunks)
    except Exception as exc:
        st.error(f"Unable to read PowerPoint: {exc}")
        return None


def preprocess_and_save(file):
    try:
        # Read the uploaded file into a DataFrame
        if file.name.endswith('.csv'):
            df = pd.read_csv(file, encoding='utf-8', na_values=['NA', 'N/A', 'missing'])
        elif file.name.endswith('.xlsx'):
            df = pd.read_excel(file, na_values=['NA', 'N/A', 'missing'])
        elif file.name.endswith('.pdf'):
            pdf_text = extract_pdf_text(file)
            if pdf_text is None:
                return None, None, None
            return None, ["document_text"], pd.DataFrame({"document_text": [pdf_text]})
        elif file.name.endswith(('.ppt', '.pptx')):
            ppt_text = extract_ppt_text(file)
            if ppt_text is None:
                return None, None, None
            return None, ["document_text"], pd.DataFrame({"document_text": [ppt_text]})
        else:
            st.error("Unsupported file format. Please upload a CSV, Excel, PDF, or PowerPoint file.")
            return None, None, None

        # Ensure string columns are properly quoted
        for col in df.select_dtypes(include=['object']):
            df[col] = df[col].astype(str).replace({r'"': '""'}, regex=True)

        # Parse dates and numeric columns
        for col in df.columns:
            if 'date' in col.lower():
                df[col] = pd.to_datetime(df[col], errors='coerce')
            elif df[col].dtype == 'object':
                try:
                    df[col] = pd.to_numeric(df[col])
                except (ValueError, TypeError):
                    # Keep as is if conversion fails
                    pass

        # Create a temporary file to save the preprocessed data
        with tempfile.NamedTemporaryFile(delete=False, suffix=".csv") as temp_file:
            temp_path = temp_file.name
            # Save the DataFrame to the temporary CSV file with quotes around string fields
            df.to_csv(temp_path, index=False, quoting=csv.QUOTE_ALL)

        return temp_path, df.columns.tolist(), df  # Return the DataFrame as well
    except Exception as e:
        st.error(f"Error processing file: {e}")
        return None, None, None


def build_document_agent(provider, api_key, document_text):
    """Build an AI agent with the full document as context for intelligent Q&A."""
    model = build_model(provider, api_key)
    word_count = len(re.findall(r"\b\w+\b", document_text or ""))

    system_prompt = f"""You are an expert document analyst AI assistant. You have been given the full text of a document to analyze.

IMPORTANT INSTRUCTIONS:
- Answer every question based ONLY on the document content provided below.
- Be thorough, specific, and cite details from the document.
- Use markdown formatting: headers, bullet points, bold text, tables where appropriate.
- If the document is a resume/CV, provide ATS analysis, strengths, gaps, and improvement suggestions.
- If it is a research paper, summarize findings, methodology, key results, and conclusions.
- If it is a presentation (PPT), summarize each slide's key points and the overall narrative.
- If it is any other document, identify its purpose, key information, and actionable insights.
- Always be helpful, detailed, and professional.

--- DOCUMENT START (approximately {word_count} words) ---
{document_text}
--- DOCUMENT END ---"""

    agent = Agent(
        model=model,
        system_message=system_prompt,
        markdown=True,
    )
    return agent


def render_visual_insights(df):
    st.subheader("📈 End-to-end business insights")
    st.caption("This section turns your data into a clear workflow story: trend, comparison, and change summaries.")

    datetime_cols = [col for col in df.columns if pd.api.types.is_datetime64_any_dtype(df[col])]
    numeric_cols = [col for col in df.columns if pd.api.types.is_numeric_dtype(df[col])]
    non_numeric_cols = [col for col in df.columns if col not in datetime_cols and col not in numeric_cols]

    if datetime_cols and numeric_cols:
        date_col = datetime_cols[0]
        value_col = next((col for col in numeric_cols if col != date_col), numeric_cols[0])

        temp_df = df[[date_col, value_col]].dropna().sort_values(date_col)
        if len(temp_df) >= 2:
            st.line_chart(temp_df.set_index(date_col)[[value_col]])

            start_value = float(temp_df.iloc[0][value_col])
            end_value = float(temp_df.iloc[-1][value_col])
            change = end_value - start_value
            if start_value != 0:
                percent_change = (change / start_value) * 100
            else:
                percent_change = None

            if change > 0:
                direction = "increasing"
            elif change < 0:
                direction = "decreasing"
            else:
                direction = "stable"

            if percent_change is not None:
                st.caption(
                    f"{value_col} is {direction} from {start_value:,.2f} to {end_value:,.2f} "
                    f"({percent_change:+.1f}%)."
                )
            else:
                st.caption(f"{value_col} is {direction} from {start_value:,.2f} to {end_value:,.2f}.")

    if non_numeric_cols and numeric_cols:
        category_col = non_numeric_cols[0]
        value_col = numeric_cols[0]
        grouped = df.groupby(category_col, dropna=False)[value_col].mean().sort_values(ascending=False).head(10)
        if not grouped.empty:
            st.subheader("📊 Category comparison")
            st.bar_chart(grouped)
            st.caption(f"Average {value_col} by {category_col}.")

            if len(grouped) > 1:
                pie_data = grouped.head(6)
                if pie_data.sum() > 0:
                    st.subheader("🥧 Share of total")
                    st.bar_chart(pie_data)

    if len(numeric_cols) >= 2:
        x_col = numeric_cols[0]
        y_col = numeric_cols[1]
        scatter_df = df[[x_col, y_col]].dropna()
        if len(scatter_df) >= 2:
            st.subheader("� Relationship view")
            st.scatter_chart(scatter_df.set_index(x_col)[[y_col]])

    if len(numeric_cols) >= 1 and len(df) > 0:
        st.subheader("🔍 Top changes")
        numeric_series = pd.to_numeric(df[numeric_cols[0]], errors='coerce').dropna()
        if len(numeric_series) >= 3:
            sorted_vals = numeric_series.sort_values()
            top_increasing = sorted_vals.tail(3)
            top_decreasing = sorted_vals.head(3)
            st.write("Highest values:")
            st.dataframe(top_increasing.to_frame(name=numeric_cols[0]))
            st.write("Lowest values:")
            st.dataframe(top_decreasing.to_frame(name=numeric_cols[0]))

    if not datetime_cols and numeric_cols and len(df) > 0:
        st.subheader("📉 Distribution")
        st.bar_chart(df[numeric_cols[0]].value_counts().head(10))
        st.caption("Quick distribution view for the selected numeric column.")


def main():
    st.set_page_config(page_title="DataLens AI", page_icon="🔍", layout="wide")
    st.markdown(
        """
        <style>
        :root {
            --bg1: #07111f;
            --bg2: #13253f;
            --panel: rgba(255,255,255,0.08);
            --border: rgba(255,255,255,0.14);
            --text: #f8fafc;
            --muted: #cbd5e1;
        }
        .stApp {
            background: radial-gradient(circle at top left, #1d4ed8 0%, transparent 28%),
                        radial-gradient(circle at bottom right, #8b5cf6 0%, transparent 25%),
                        linear-gradient(135deg, var(--bg1), var(--bg2));
            color: var(--text);
        }
        .block-container {
            padding-top: 1.2rem;
            padding-bottom: 2rem;
        }
        div[data-testid="stSidebar"] {
            background: linear-gradient(180deg, #0b1220 0%, #111c2f 100%);
            border-right: 1px solid var(--border);
            box-shadow: 6px 0 24px rgba(0,0,0,0.35);
        }
        div[data-testid="stSidebar"] .stTextInput > label,
        div[data-testid="stSidebar"] .stSelectbox > label,
        div[data-testid="stSidebar"] .stMarkdown {
            color: var(--text);
        }
        .stAlert, .stSuccess, .stWarning, .stInfo {
            border-radius: 14px;
            box-shadow: 0 10px 25px rgba(0,0,0,0.18);
        }
        div[data-testid="stVerticalBlock"] > div > div {
            border-radius: 18px;
            border: 1px solid var(--border);
            background: var(--panel);
            backdrop-filter: blur(12px);
            box-shadow: 0 12px 30px rgba(0,0,0,0.24);
            padding: 16px;
            margin-bottom: 14px;
        }
        .stDataFrame, .stTable {
            border-radius: 14px;
            overflow: hidden;
            box-shadow: inset 0 1px 0 rgba(255,255,255,0.08);
        }
        .stButton > button {
            border-radius: 999px;
            background: linear-gradient(90deg, #3b82f6, #8b5cf6);
            color: white;
            border: none;
            box-shadow: 0 10px 20px rgba(59,130,246,0.25);
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.title("🔍 DataLens AI")
    st.caption("Upload any file — CSV, Excel, PDF, or PowerPoint — and let AI uncover insights, trends, and answers instantly.")

    # Sidebar for API keys
    with st.sidebar:
        st.header("⚙️ Model Settings")
        st.caption("Choose the provider for your assistant.")
        provider = st.selectbox("Choose model provider", ["google", "openai"], index=0)
        openai_key = resolve_openai_key(st.session_state)

        if provider == "openai":
            if openai_key:
                st.success("OpenAI key loaded automatically.")
            else:
                openai_key = st.text_input("Enter your OpenAI API key:", type="password")
                if openai_key:
                    st.session_state.openai_key = openai_key
                    st.success("OpenAI key saved!")
                else:
                    st.warning("Please enter your OpenAI API key to proceed.")
        else:
            google_key = os.getenv("GOOGLE_API_KEY") or st.text_input("Enter your Google API key:", type="password")
            if google_key:
                st.session_state.openai_key = google_key
                st.success("Google key ready.")
            else:
                st.warning("Please enter your Google API key to proceed.")

        st.divider()
        st.caption("💡 Tip: CSV/Excel → charts & SQL analysis · PDF/PPT → AI summary & chat")

    # File upload widget
    st.subheader("📁 Upload your data")
    uploaded_file = st.file_uploader("Upload a CSV, Excel, PDF, or PowerPoint file", type=["csv", "xlsx", "pdf", "ppt", "pptx"], label_visibility="visible")

    if uploaded_file is not None:
        # Preprocess and save the uploaded file
        temp_path, columns, df = preprocess_and_save(uploaded_file)

        if columns and df is not None:
            st.success("File uploaded successfully.")

            # --- API key check (required for BOTH documents and data) ---
            if not openai_key:
                st.warning("Please enter a valid API key for the selected provider to proceed.")
                return

            if columns == ["document_text"]:
                document_text = df.iloc[0, 0]
                word_count = len(re.findall(r"\b\w+\b", document_text or ""))

                # --- Document info header ---
                st.subheader("📄 Document uploaded")
                col1, col2 = st.columns(2)
                col1.metric("📝 Word count", f"{word_count:,}")
                col2.metric("📄 File type", uploaded_file.name.split('.')[-1].upper())

                with st.expander("📄 View extracted document text", expanded=False):
                    st.text_area("Extracted text", value=document_text, height=300, disabled=True)

                # --- Build AI document agent ---
                doc_agent = build_document_agent(provider, openai_key, document_text)

                # --- Auto-generate AI summary on first upload ---
                file_key = f"doc_summary_{uploaded_file.name}_{uploaded_file.size}"
                if file_key not in st.session_state:
                    with st.spinner("🤖 AI is analyzing your document..."):
                        try:
                            summary_prompt = (
                                "Analyze this document comprehensively. Provide:\n"
                                "1. **Document Type** — What kind of document is this?\n"
                                "2. **Executive Summary** — A clear, concise overview\n"
                                "3. **Key Findings / Main Points** — The most important information (use bullet points)\n"
                                "4. **Important Details** — Statistics, dates, names, or data worth noting\n"
                                "5. **Actionable Insights / Recommendations** — What should the reader do with this information?\n\n"
                                "Be thorough and specific. Use markdown formatting."
                            )
                            response = doc_agent.run(summary_prompt)
                            summary_content = response.content if hasattr(response, 'content') else str(response)
                            st.session_state[file_key] = summary_content
                        except Exception as e:
                            st.session_state[file_key] = f"⚠️ Could not generate summary: {e}"

                # --- Display AI summary ---
                st.subheader("🤖 AI Analysis")
                st.markdown(st.session_state.get(file_key, ""))

                st.divider()

                # --- Chat interface ---
                st.subheader("💬 Chat with your document")
                st.caption("Ask any question — the AI has read your entire document.")

                # Initialize chat history for this file
                chat_key = f"doc_chat_{uploaded_file.name}_{uploaded_file.size}"
                if chat_key not in st.session_state:
                    st.session_state[chat_key] = []

                # Display chat history
                for msg in st.session_state[chat_key]:
                    with st.chat_message(msg["role"]):
                        st.markdown(msg["content"])

                # Chat input
                user_query = st.chat_input("Ask about your document...")
                if user_query:
                    # Show user message
                    st.session_state[chat_key].append({"role": "user", "content": user_query})
                    with st.chat_message("user"):
                        st.markdown(user_query)

                    # Get AI response
                    with st.chat_message("assistant"):
                        with st.spinner("Thinking..."):
                            try:
                                # Build context from chat history
                                history_context = ""
                                for msg in st.session_state[chat_key][:-1]:
                                    role_label = "User" if msg["role"] == "user" else "Assistant"
                                    history_context += f"{role_label}: {msg['content']}\n\n"

                                full_query = user_query
                                if history_context:
                                    full_query = (
                                        f"Previous conversation:\n{history_context}\n"
                                        f"Current question: {user_query}"
                                    )

                                response = doc_agent.run(full_query)
                                response_content = response.content if hasattr(response, 'content') else str(response)
                            except Exception as e:
                                response_content = f"⚠️ Error: {e}\n\nPlease try rephrasing your question."

                        st.markdown(response_content)
                        st.session_state[chat_key].append({"role": "assistant", "content": response_content})

            else:

                render_visual_insights(df)
                st.subheader("🧾 Uploaded data")
                st.dataframe(df, use_container_width=True)
                st.subheader("🧠 Data preview")
                st.write("Columns:", columns)

                # Initialize DuckDbTools
                duckdb_tools = DuckDbTools()

                # Load the CSV file into DuckDB as a table
                duckdb_tools.load_local_csv_to_table(
                    path=temp_path,
                    table="uploaded_data",
                )

                # Initialize the Agent with DuckDB and Pandas tools
                model = build_model(provider, openai_key)
                data_analyst_agent = Agent(
                    model=model,
                    tools=[duckdb_tools, PandasTools()],
                    system_message="You are an expert data analyst. Use the 'uploaded_data' table to answer user queries. Generate SQL queries using DuckDB tools to solve the user's query. Provide clear and concise answers with the results.",
                    markdown=True,
                )

                # Initialize code storage in session state
                if "generated_code" not in st.session_state:
                    st.session_state.generated_code = None

                st.subheader("💬 Ask your data a question")
                user_query = st.text_area(
                    "Ask a question about your data:",
                    placeholder="Example: Show the trend over time and explain the biggest increase.",
                )
                st.info("The assistant will analyze your uploaded file and summarize the key insights.")

                if st.button("Submit Query"):
                    if user_query.strip() == "":
                        st.warning("Please enter a query.")
                    else:
                        try:
                            # Show loading spinner while processing
                            with st.spinner('Processing your query...'):
                                # Get the response from the agent
                                response = data_analyst_agent.run(user_query)

                                # Extract the content from the response object
                                if hasattr(response, 'content'):
                                    response_content = response.content
                                else:
                                    response_content = str(response)

                            # Display the response in Streamlit
                            st.markdown(response_content)

                        except Exception as e:
                            st.error(f"Error generating response from the agent: {e}")
                            st.error("Please try rephrasing your query or check if the data format is correct.")


if __name__ == "__main__":
    main()